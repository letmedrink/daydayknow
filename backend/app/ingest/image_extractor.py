"""从 PDF/PPTX/DOCX 提取嵌入图片，保存到 wiki/media/ 目录。"""

import hashlib
import io
from pathlib import Path
from typing import Optional


def extract_images(filename: str, content: bytes, media_dir: Path) -> list[dict]:
    """提取文件中的嵌入图片。

    Args:
        filename: 原始文件名
        content: 文件二进制内容
        media_dir: 图片保存目录 (wiki/media/<slug>/)

    Returns:
        图片信息列表 [{filename, rel_path, page, sha256, width, height}]
    """
    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        return _extract_from_pdf(content, media_dir)
    elif ext == ".pptx":
        return _extract_from_pptx(content, media_dir)
    elif ext == ".docx":
        return _extract_from_docx(content, media_dir)
    return []


def _extract_from_pdf(content: bytes, media_dir: Path) -> list[dict]:
    """从 PDF 提取图片。"""
    try:
        import fitz
    except ImportError:
        return []

    images = []
    doc = fitz.open(stream=content, filetype="pdf")

    for page_idx in range(len(doc)):
        page = doc[page_idx]
        image_list = page.get_images(full=True)

        for img_idx, img_info in enumerate(image_list):
            xref = img_info[0]
            try:
                base_image = doc.extract_image(xref)
                if not base_image:
                    continue

                img_bytes = base_image["image"]
                img_ext = base_image.get("ext", "png")
                sha256 = hashlib.sha256(img_bytes).hexdigest()

                img_filename = f"page{page_idx + 1}_img{img_idx + 1}.{img_ext}"
                media_dir.mkdir(parents=True, exist_ok=True)
                img_path = media_dir / img_filename
                img_path.write_bytes(img_bytes)

                images.append({
                    "filename": img_filename,
                    "rel_path": f"media/{media_dir.name}/{img_filename}",
                    "page": page_idx + 1,
                    "sha256": sha256,
                    "width": base_image.get("width", 0),
                    "height": base_image.get("height", 0),
                })
            except Exception:
                continue

    doc.close()
    return images


def _extract_from_pptx(content: bytes, media_dir: Path) -> list[dict]:
    """从 PPTX 提取图片。"""
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        return []

    images = []
    prs = Presentation(io.BytesIO(content))

    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image = shape.image
                    img_bytes = image.blob
                    ext = image.content_type.split("/")[-1]
                    if ext == "jpeg":
                        ext = "jpg"
                    sha256 = hashlib.sha256(img_bytes).hexdigest()

                    img_filename = f"slide{slide_idx + 1}_img{shape_idx + 1}.{ext}"
                    media_dir.mkdir(parents=True, exist_ok=True)
                    img_path = media_dir / img_filename
                    img_path.write_bytes(img_bytes)

                    images.append({
                        "filename": img_filename,
                        "rel_path": f"media/{media_dir.name}/{img_filename}",
                        "page": slide_idx + 1,
                        "sha256": sha256,
                        "width": 0,
                        "height": 0,
                    })
                except Exception:
                    continue

    return images


def _extract_from_docx(content: bytes, media_dir: Path) -> list[dict]:
    """从 DOCX 提取图片。"""
    try:
        from docx import Document
    except ImportError:
        return []

    images = []
    doc = Document(io.BytesIO(content))

    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                image_part = rel.target_part
                img_bytes = image_part.blob
                ext = Path(image_part.partname).suffix.lstrip(".")
                if not ext:
                    ext = "png"
                sha256 = hashlib.sha256(img_bytes).hexdigest()

                img_filename = f"img_{sha256[:8]}.{ext}"
                media_dir.mkdir(parents=True, exist_ok=True)
                img_path = media_dir / img_filename
                img_path.write_bytes(img_bytes)

                images.append({
                    "filename": img_filename,
                    "rel_path": f"media/{media_dir.name}/{img_filename}",
                    "page": 0,
                    "sha256": sha256,
                    "width": 0,
                    "height": 0,
                })
            except Exception:
                continue

    return images
