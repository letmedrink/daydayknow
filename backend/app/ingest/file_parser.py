"""文档解析 — 从 PDF/PPTX/DOCX/TXT/MD 提取文本内容。"""

import io
from pathlib import Path
from typing import Optional


def parse_file(filename: str, content: bytes) -> str:
    """解析文件内容，返回纯文本。

    支持格式: .pdf, .pptx, .docx, .txt, .md, .csv, .json
    """
    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        return _parse_pdf(content)
    elif ext == ".pptx":
        return _parse_pptx(content)
    elif ext == ".docx":
        return _parse_docx(content)
    elif ext in (".txt", ".md", ".csv", ".json", ".py", ".js", ".ts", ".yaml", ".yml", ".toml"):
        return content.decode("utf-8", errors="replace")
    else:
        # 尝试作为文本读取
        try:
            return content.decode("utf-8")
        except UnicodeDecodeError:
            return f"[不支持的文件格式: {ext}]"


def _parse_pdf(content: bytes) -> str:
    """解析 PDF 文件。"""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=content, filetype="pdf")
        pages = []
        for i, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                pages.append(f"--- 第 {i + 1} 页 ---\n{text}")
        doc.close()
        return "\n\n".join(pages) if pages else "[PDF 文件无可提取文本]"
    except ImportError:
        return "[需要安装 PyMuPDF: pip install pymupdf]"
    except Exception as e:
        return f"[PDF 解析失败: {e}]"


def _parse_pptx(content: bytes) -> str:
    """解析 PowerPoint 文件。"""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slides.append(f"--- 幻灯片 {i + 1} ---\n" + "\n".join(texts))
        return "\n\n".join(slides) if slides else "[PPTX 文件无可提取文本]"
    except ImportError:
        return "[需要安装 python-pptx: pip install python-pptx]"
    except Exception as e:
        return f"[PPTX 解析失败: {e}]"


def _parse_docx(content: bytes) -> str:
    """解析 Word 文件。"""
    try:
        from docx import Document
        doc = Document(io.BytesIO(content))
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)
        # 也提取表格内容
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    paragraphs.append(row_text)
        return "\n\n".join(paragraphs) if paragraphs else "[DOCX 文件无可提取文本]"
    except ImportError:
        return "[需要安装 python-docx: pip install python-docx]"
    except Exception as e:
        return f"[DOCX 解析失败: {e}]"
